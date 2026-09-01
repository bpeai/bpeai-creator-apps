from __future__ import annotations

from pathlib import Path

import pytest

from pptx import Presentation

from bpeai_creator_sdk.artifacts import (
    build_evaluation_pptx,
    build_slide_pack_from_evaluation,
)
from bpeai_creator_sdk.local_format import format_result_text
from bpeai_creator_sdk.local_run import is_selector_result
from bpeai_creator_sdk.output import validate_output
from bpeai_creator_sdk.sme import (
    missing_report_headings,
    validate_dir_code,
)

from pack_paths import load_platform_pack_or_skip


@pytest.fixture(scope="module")
def py_root() -> Path:
    return Path(__file__).resolve().parents[3]


@pytest.fixture(scope="module")
def mixing_pack(py_root: Path):
    return load_platform_pack_or_skip("mixing", py_root)


def test_pack_loads_outlines(mixing_pack):
    assert mixing_pack.required_report_headings()
    assert "Validated DIR" in mixing_pack.required_report_headings()
    assert mixing_pack.pptx_outline.get("slide_count") == 7
    assert mixing_pack.fragment("workflow")
    assert mixing_pack.common_code_entries("media_preparation")[0]["caption"]


def test_missing_report_headings():
    md = "# Validated DIR\n# Design basis\n"
    missing = missing_report_headings(
        md,
        ["Validated DIR", "Design basis", "Strong-fit mixing types"],
    )
    assert missing == ["Strong-fit mixing types"]


def test_dir_code_still_validates(mixing_pack):
    ok = validate_dir_code(mixing_pack, "media_preparation", "2-1-2-3-1-1")
    assert ok.ok


def test_build_evaluation_pptx_smoke(tmp_path: Path, mixing_pack):
    fixture = {
        "system_name": "Media Prep Vessel",
        "dir_code": "2-1-2-3-1-1",
        "selected_model": "Top-entry low-shear axial hydrofoil agitator",
        "recommended_basis": "Top-entry low-shear axial hydrofoil agitator",
        "alternate_basis": "Aseptic magnetic bottom mixer",
        "design_basis": "Readily soluble dry powder, manual top-charge, CIP/SIP stainless.",
        "failure_modes": ["Clumping", "Foam", "Seal SIP issue", "Low-fill dead zone"],
        "objectives": ["Dissolve", "Homogenize", "Hygienic"],
        "mixing_options": [
            {
                "name": "Top-entry low-shear axial hydrofoil agitator",
                "fit": "best",
                "pros": ["low foam", "CIP/SIP", "turndown"],
                "cons": ["seal care"],
                "manufacturers": ["Lightnin A310"],
            }
        ],
        "evaluation_matrix": [
            {
                "option": "Hydrofoil",
                "technical_fit": "Best",
                "gmp": "High",
                "scale_up_risk": "Low",
                "cost_schedule": "Best",
                "reliability": "High",
                "rank": 1,
            }
        ],
        "preliminary_specs": ["316L", "VFD", "Baffles"],
        "manufacturers": ["Lightnin", "Alfa Laval"],
        "do_not_specify": ["Rushton-only"],
        "rationale": "Best fit for manual charge media prep.",
        "datasheet_markdown": "# Design basis\nTest body.\n\n# Option evaluation\nDetails here.\n",
    }
    out = tmp_path / "media_prep.pptx"
    slide_pack = build_slide_pack_from_evaluation(fixture)
    path = build_evaluation_pptx(
        fixture,
        outline=mixing_pack.pptx_outline,
        output_path=out,
        slide_pack=slide_pack,
    )
    assert path.is_file()
    prs = Presentation(str(path))
    assert len(prs.slides) == 7
    assert abs(prs.slide_width.inches - 13.333) < 0.01


def test_slide3_objective_cards_keep_long_titles(tmp_path: Path):
    """Slide 3 titles must wrap in the card, not hard-clip at 22 characters."""
    fixture = {
        "system_name": "MF Harvest Clarification Skid",
        "dir_code": "3-3-3-3-3-3-3",
        "selected_model": "Closed MF-TFF harvest skid",
        "recommended_basis": "Closed MF-TFF harvest skid",
        "objectives": ["Clarify harvest", "Operate at target flux"],
        "failure_modes": ["TMP runaway", "Turbidity breakthrough"],
    }
    slide_pack = build_slide_pack_from_evaluation(fixture)
    slide_pack["slides"][2]["process_steps"] = [
        {"n": 1, "title": "Clarify harvest", "detail": "5,000–10,000 L in ≤8 h; ≥95% recovery."},
        {"n": 2, "title": "Operate at target flux", "detail": "100–155 LMH via constant-flux with TMP limit."},
        {
            "n": 3,
            "title": "Manage hydraulics",
            "detail": "Crossflow 300–600 L/m²/h; TMP 0.8–2.0 bar controlled.",
        },
        {
            "n": 4,
            "title": "Maintain thermal consistency",
            "detail": "≤2 °C temperature rise across skid at design conditions.",
        },
    ]
    path = build_evaluation_pptx(
        fixture,
        outline=None,
        output_path=tmp_path / "slide3.pptx",
        slide_pack=slide_pack,
    )
    prs = Presentation(str(path))
    texts = []
    for shape in prs.slides[2].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    joined = "\n".join(texts)
    assert "Maintain thermal consistency" in joined
    assert "cons…" not in joined
    assert "design conditions" in joined
    assert "bar controlled" in joined


def test_build_evaluation_pdf_and_reference_decks(tmp_path: Path, mixing_pack):
    from bpeai_creator_sdk import build_evaluation_pdf, list_reference_decks

    pdf = build_evaluation_pdf(
        {
            "system_name": "Media Prep Vessel",
            "dir_code": "2-1-2-3-1-1",
            "datasheet_markdown": "# Design basis\n\nSelection implication narrative.\n\n## Option evaluation\n\n- Hydrofoil best fit\n",
            "selected_model": "Hydrofoil",
        },
        output_path=tmp_path / "eval.pdf",
    )
    assert pdf.is_file()
    assert pdf.stat().st_size > 500

    decks = list_reference_decks(mixing_pack.path, outline=mixing_pack.pptx_outline)
    assert decks
    assert any(d["name"].endswith(".pptx") for d in decks)


def test_format_evaluation_mentions_pptx_prompt():
    text = format_result_text(
        {
            "phase": "evaluation",
            "schema_version": "equipment_selector_v1",
            "system_name": "Media Prep",
            "dir_code": "2-1-2-3-1-1",
            "selected_model": "Hydrofoil",
            "recommended_basis": "Hydrofoil",
            "rationale": "Fit",
            "mixing_options": [{"name": "Hydrofoil", "fit": "best", "pros": ["low foam"]}],
            "pptx_prompt": "Would you like a presentation-ready PPTX file? Reply pptx or y.",
        }
    )
    assert "Recommended basis" in text
    assert "pptx" in text.lower()


def test_validate_preserves_session_extras_when_merged():
    """run_agent merges validated fields back onto extras used for pptx session."""
    raw = {
        "schema_version": "equipment_selector_v1",
        "equipment_tag": "AG-101",
        "selected_model": "Hydrofoil",
        "equipment_system": "mixing",
        "rationale": "Fit",
        "creator_attribution": {"display_name": "T", "app_id": "equipment_evaluator"},
        "phase": "evaluation",
        "system_name": "Media Prep Vessel",
        "application": "biopharmaceutical",
        "dir_code": "3-1-3-3-3-4",
        "pptx_prompt": "Reply pptx or y.",
        "artifacts": {"markdown_path": "artifacts/x.md"},
    }
    assert is_selector_result(raw)
    validated = validate_output(raw).model_dump()
    # Bare validate strips extras (the bug that broke pptx session storage).
    assert "phase" not in validated or validated.get("phase") is None
    merged = dict(raw)
    merged.update(validated)
    assert merged["phase"] == "evaluation"
    assert merged["dir_code"] == "3-1-3-3-3-4"
    assert merged["system_name"] == "Media Prep Vessel"
    assert merged["selected_model"] == "Hydrofoil"
