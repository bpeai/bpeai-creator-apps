"""Build 7-slide mixing evaluation decks matching the Life Science Mixing Systems Expert style.

Visual system derived from:
``py/knowledge/mixing/references/chromatography_resin_slurry_tank_agitator_evaluation.pptx``
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Mapping

# Palette (reference deck)
NAVY = "17324D"
TEAL = "00A398"
BLUE = "2962A3"
GRAY = "6B7280"
BODY = "1F2937"
WHITE = "FFFFFF"
CARD_FILL = "F3F6F9"
CHIP_BLUE_FILL = "E8F1FB"
CHIP_TEAL_FILL = "E6F7F6"
PANEL_RIGHT = "F7FAFC"

FONT_DISPLAY = "Aptos Display"
FONT_BODY = "Aptos"


def _rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    h = hex_color.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _as_list(value: Any) -> List[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [str(v).strip() for v in value if str(v).strip()]
    text = str(value).strip()
    return [text] if text else []


def _truncate(text: str, limit: int) -> str:
    text = " ".join(str(text or "").split())
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "…"


def _emu_to_inches(emu: int) -> float:
    return float(emu) / 914400.0


def _fit_font_pt(
    text: str,
    *,
    width_emu: int,
    height_emu: int,
    preferred_pt: float,
    min_pt: float = 9.0,
    paragraphs: int = 1,
) -> float:
    """Estimate a font size that keeps wrapped text inside a fixed box."""
    clean = " ".join(str(text or "").split())
    if not clean:
        return preferred_pt
    size = float(preferred_pt)
    para_count = max(1, paragraphs)
    while size > min_pt:
        width_in = max(0.2, _emu_to_inches(width_emu))
        height_in = max(0.15, _emu_to_inches(height_emu))
        # Aptos-ish average glyph width ≈ 0.52em
        chars_per_line = max(6, int(width_in * 72.0 / (size * 0.52)))
        line_height_in = (size / 72.0) * 1.2
        max_lines = max(1, int(height_in / line_height_in))
        # Reserve one line per paragraph, then remaining for wrap
        if para_count > max_lines:
            size -= 0.5
            continue
        wrap_budget = max_lines * chars_per_line
        if len(clean) <= wrap_budget:
            return size
        size -= 0.5
    return min_pt


def _set_run(paragraph, text: str, *, size_pt: float, bold: bool | None, color: str, font_name: str) -> None:
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.size = __import__("pptx.util", fromlist=["Pt"]).Pt(size_pt)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = _rgb(color)


def _fill_textbox(
    shape,
    items: List[str] | str,
    *,
    width_emu: int,
    height_emu: int,
    preferred_pt: float,
    min_pt: float = 9.0,
    bold: bool | None = False,
    color: str = BODY,
    font_name: str = FONT_BODY,
    center: bool = False,
    bullet: bool = False,
    max_items: int | None = None,
    char_cap: int | None = None,
) -> None:
    """Write one or more paragraphs into a text box, shrinking font to fit."""
    from pptx.enum.text import PP_ALIGN

    if isinstance(items, str):
        lines = [items] if items.strip() else []
    else:
        lines = [str(x).strip() for x in items if str(x).strip()]
    if max_items is not None:
        lines = lines[:max_items]
    if char_cap is not None:
        lines = [_truncate(x, char_cap) for x in lines]
    if not lines:
        lines = [""]

    joined = " ".join(lines)
    size = _fit_font_pt(
        joined,
        width_emu=width_emu,
        height_emu=height_emu,
        preferred_pt=preferred_pt,
        min_pt=min_pt,
        paragraphs=len(lines),
    )
    # Extra shrink when many bullets crowd the box
    if len(lines) >= 4 and size > min_pt:
        size = max(min_pt, size - 1.0)
    if len(lines) >= 6 and size > min_pt:
        size = max(min_pt, size - 1.0)

    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if center:
            p.alignment = PP_ALIGN.CENTER
        display = f"• {line}" if bullet and line else line
        _set_run(p, display, size_pt=size, bold=bold, color=color, font_name=font_name)


def _add_textbox(slide, left, top, width, height):
    from pptx.util import Emu

    return slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))


def _add_rect(slide, left, top, width, height, fill_hex: str, *, line: bool = False):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if not line:
        shape.line.fill.background()
    return shape


def _add_line(slide, left, top, width):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(19050))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb("E5E7EB")
    shape.line.fill.background()
    return shape


def _footer(slide, *, slide_no: int, dir_code: str) -> None:
    from pptx.enum.text import PP_ALIGN

    left = _add_textbox(slide, 502920, 6510528, 3657600, 182880)
    p = left.text_frame.paragraphs[0]
    _set_run(p, f"Project-team summary • Slide {slide_no}", size_pt=8.5, bold=False, color=GRAY, font_name=FONT_BODY)
    right = _add_textbox(slide, 8046720, 6510528, 3657600, 182880)
    rp = right.text_frame.paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    _set_run(
        rp,
        f"Basis: user DIR code {dir_code}" if dir_code else "Basis: validated DIR",
        size_pt=8.5,
        bold=False,
        color=GRAY,
        font_name=FONT_BODY,
    )


def _eyebrow_and_title(slide, eyebrow: str, title: str) -> None:
    eb = _add_textbox(slide, 502920, 228600, 10698480, 219456)
    _set_run(eb.text_frame.paragraphs[0], eyebrow, size_pt=9.5, bold=True, color=GRAY, font_name=FONT_BODY)
    hd = _add_textbox(slide, 502920, 502920, 10698480, 438912)
    _set_run(hd.text_frame.paragraphs[0], title, size_pt=28, bold=True, color=NAVY, font_name=FONT_DISPLAY)
    _add_line(slide, 502920, 1060704, 11155680)


def default_reference_path(pack_path: Path | None = None) -> Path | None:
    candidates = []
    if pack_path is not None:
        candidates.append(Path(pack_path) / "references" / "chromatography_resin_slurry_tank_agitator_evaluation.pptx")
        candidates.append(Path(pack_path) / "references" / "media_preparation_vessel_mixing_evaluation.pptx")
    # relative fallbacks
    here = Path(__file__).resolve()
    for root in here.parents:
        knowledge = root / "knowledge" / "mixing" / "references"
        candidates.append(knowledge / "chromatography_resin_slurry_tank_agitator_evaluation.pptx")
        candidates.append(knowledge / "media_preparation_vessel_mixing_evaluation.pptx")
    for c in candidates:
        if c.is_file():
            return c
    return None


def build_slide_pack_from_evaluation(result: Mapping[str, Any]) -> Dict[str, Any]:
    """Deterministic fallback slide pack when LLM is unavailable."""
    system = str(result.get("system_name") or result.get("equipment_name") or "Mixing system")
    dir_code = str(result.get("dir_code") or "")
    recommended = str(result.get("recommended_basis") or result.get("selected_model") or "")
    alternate = str(result.get("alternate_basis") or "")
    title_lines = system.split()
    if len(title_lines) >= 2:
        mid = max(1, len(title_lines) // 2)
        title_lines = [" ".join(title_lines[:mid]), " ".join(title_lines[mid:])]
    else:
        title_lines = [system]

    decoded = result.get("decoded_dir") or []
    cards = []
    for row in decoded[:6]:
        if isinstance(row, Mapping):
            cards.append(
                {
                    "label": str(row.get("label") or "").upper(),
                    "value": str(row.get("option_text") or ""),
                    "accent": "objective" in str(row.get("label") or "").lower()
                    or "duty" in str(row.get("label") or "").lower(),
                }
            )
    if not cards:
        for spec in (result.get("key_specs") or [])[:6]:
            if isinstance(spec, Mapping):
                cards.append(
                    {
                        "label": str(spec.get("key") or "").upper(),
                        "value": str(spec.get("value") or ""),
                        "accent": False,
                    }
                )

    options = [o for o in (result.get("mixing_options") or []) if isinstance(o, Mapping)]
    matrix = [m for m in (result.get("evaluation_matrix") or []) if isinstance(m, Mapping)]
    if not matrix:
        for i, opt in enumerate(options[:6], start=1):
            matrix.append(
                {
                    "option": opt.get("name") or f"Option {i}",
                    "technical_fit": str(opt.get("fit") or "").title(),
                    "gmp": "High",
                    "scale_up_risk": "Medium",
                    "cost_schedule": "Medium",
                    "reliability": "High",
                    "rank": i,
                }
            )

    process_steps = []
    for i, obj in enumerate(_as_list(result.get("objectives"))[:4], start=1):
        process_steps.append({"n": i, "title": _truncate(obj, 28), "detail": _truncate(obj, 40)})
    if not process_steps:
        process_steps = [
            {"n": 1, "title": "Charge", "detail": "Controlled liquid charge"},
            {"n": 2, "title": "Mix / dissolve", "detail": "Achieve homogeneity"},
            {"n": 3, "title": "Sample", "detail": "Confirm quality attributes"},
            {"n": 4, "title": "Transfer", "detail": "Release homogeneous batch"},
        ]

    return {
        "slides": [
            {
                "id": "title",
                "title_lines": title_lines[:3],
                "subtitle": _truncate(
                    result.get("dir_summary")
                    or result.get("design_basis")
                    or f"Mixing technology evaluation for {system}.",
                    160,
                ),
                "dir_badge": f"Validated DIR: {dir_code}" if dir_code else "Validated DIR",
                "summary_badge": "Project-team summary",
                "hero_tags": _as_list(result.get("objectives"))[:3]
                or ["technical fit", "GMP ready", "vendor available"],
                "hero_headline": [
                    _truncate(recommended, 34) or "Recommended basis",
                    "as recommended basis",
                    "for this DIR",
                ],
            },
            {
                "id": "design_basis",
                "eyebrow": f"Agitator Selection / {system}",
                "heading": "Design basis from DIR code",
                "cards": cards[:6],
                "selection_implication": _truncate(
                    result.get("design_basis") or result.get("rationale") or "",
                    220,
                ),
            },
            {
                "id": "objectives",
                "eyebrow": f"Agitator Selection / {system}",
                "heading": "Mixing objectives, constraints and failure modes",
                "process_steps": process_steps,
                "failure_modes": _as_list(result.get("failure_modes"))[:4]
                or ["Poor wetting / floating solids", "Foam and air entrainment", "Dead zones delay uniformity"],
                "target_outcome": _truncate(
                    result.get("recommended_basis")
                    or "A repeatable, GMP-suitable batch with confirmed homogeneity before transfer.",
                    160,
                ),
            },
            {
                "id": "options",
                "eyebrow": f"Agitator Selection / {system}",
                "heading": "Realistic mixing-system options",
                "rows": [
                    {
                        "name": str(o.get("name") or ""),
                        "fit": str(o.get("fit") or ""),
                        "notes": _truncate(
                            "; ".join(_as_list(o.get("industrial_applications"))[:2])
                            or "; ".join(_as_list(o.get("pros"))[:2]),
                            48,
                        ),
                    }
                    for o in options[:6]
                ],
                "recommendation_line": _truncate(
                    f"Recommendation: proceed with {recommended}"
                    + (f"; keep {alternate} as alternate." if alternate else "."),
                    180,
                ),
            },
            {
                "id": "matrix",
                "eyebrow": f"Agitator Selection / {system}",
                "heading": "Option evaluation matrix",
                "rows": matrix[:6],
                "decision_logic": _truncate(result.get("rationale") or "", 200),
            },
            {
                "id": "recommendation",
                "eyebrow": f"Agitator Selection / {system}",
                "heading": "Recommended basis and alternate option",
                "recommended": recommended,
                "recommended_why": _as_list(result.get("objectives"))[:3]
                or _truncate(result.get("rationale") or "", 120).split(". ")[:3],
                "pros": _as_list((options[0] if options else {}).get("pros"))[:4],
                "cons": _as_list((options[0] if options else {}).get("cons"))[:3],
                "alternate": alternate,
                "alternate_note": _truncate(
                    (options[1].get("name") if len(options) > 1 else "")
                    and f"Use when procurement or platform constraints favor {options[1].get('name')}.",
                    140,
                ),
            },
            {
                "id": "specs",
                "eyebrow": f"Agitator Selection / {system}",
                "heading": "Preliminary specification points / vendors / references",
                "specs": _as_list(result.get("preliminary_specs"))[:7],
                "manufacturers": _as_list(result.get("manufacturers"))[:8],
                "do_not_specify": _as_list(result.get("do_not_specify"))[:5],
                "references": _as_list(result.get("source_basis"))[:5],
            },
        ],
        "dir_code": dir_code,
        "system_name": system,
    }


def build_evaluation_pptx(
    result: Mapping[str, Any],
    *,
    outline: Mapping[str, Any] | None = None,
    output_path: Path | str,
    slide_pack: Mapping[str, Any] | None = None,
    template_path: Path | str | None = None,
    pack_path: Path | str | None = None,
) -> Path:
    """Build a styled 7-slide evaluation deck.

    Prefer ``slide_pack`` from the LLM. Falls back to deterministic packing from ``result``.
    ``template_path`` is retained for compatibility (style is recreated to match reference).
    """
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    _ = outline, template_path  # style is coded from the reference geometry

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    pack = dict(slide_pack) if isinstance(slide_pack, Mapping) else build_slide_pack_from_evaluation(result)
    slides = pack.get("slides") or []
    if not isinstance(slides, list) or len(slides) < 7:
        # merge/pad with fallback
        fallback = build_slide_pack_from_evaluation(result)
        fb_slides = fallback.get("slides") or []
        merged = []
        for i in range(7):
            if i < len(slides) and isinstance(slides[i], Mapping):
                merged.append(slides[i])
            else:
                merged.append(fb_slides[i])
        slides = merged
        pack["dir_code"] = pack.get("dir_code") or fallback.get("dir_code")
        pack["system_name"] = pack.get("system_name") or fallback.get("system_name")

    dir_code = str(pack.get("dir_code") or result.get("dir_code") or "")
    system = str(pack.get("system_name") or result.get("system_name") or "Mixing system")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    t1 = slides[0] if isinstance(slides[0], Mapping) else {}
    # right panel
    _add_rect(s1, 7589520, 0, 4599432, 6858000, PANEL_RIGHT, line=False)
    title_box = _add_textbox(s1, 658368, 1207008, 6217920, 1400000)
    lines = _as_list(t1.get("title_lines")) or system.split()
    _fill_textbox(
        title_box,
        [_truncate(x, 28) for x in lines[:3]],
        width_emu=6217920,
        height_emu=1400000,
        preferred_pt=34,
        min_pt=22,
        bold=True,
        color=NAVY,
        font_name=FONT_DISPLAY,
    )
    sub = _add_textbox(s1, 694944, 2743200, 5852160, 800000)
    _fill_textbox(
        sub,
        _truncate(t1.get("subtitle") or "", 140),
        width_emu=5852160,
        height_emu=800000,
        preferred_pt=16,
        min_pt=12,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
    )
    # chips
    _add_rect(s1, 694944, 3794760, 2514600, 310896, CHIP_BLUE_FILL)
    chip1 = _add_textbox(s1, 749808, 3849624, 2404872, 200000)
    _fill_textbox(
        chip1,
        _truncate(t1.get("dir_badge") or f"Validated DIR: {dir_code}", 32),
        width_emu=2404872,
        height_emu=200000,
        preferred_pt=10,
        min_pt=8,
        bold=True,
        color=BLUE,
        font_name=FONT_BODY,
        center=True,
    )
    _add_rect(s1, 3401568, 3794760, 2176272, 310896, CHIP_TEAL_FILL)
    chip2 = _add_textbox(s1, 3456432, 3849624, 2066544, 200000)
    _fill_textbox(
        chip2,
        _truncate(str(t1.get("summary_badge") or "Project-team summary"), 28),
        width_emu=2066544,
        height_emu=200000,
        preferred_pt=10,
        min_pt=8,
        bold=True,
        color=TEAL,
        font_name=FONT_BODY,
        center=True,
    )
    # hero tags + headline on right
    tags_box = _add_textbox(s1, 8796528, 2212848, 2000000, 1200000)
    tags = _as_list(t1.get("hero_tags"))[:3] or ["fit", "GMP", "scale-up"]
    _fill_textbox(
        tags_box,
        [_truncate(tag, 18) for tag in tags],
        width_emu=2000000,
        height_emu=1200000,
        preferred_pt=13,
        min_pt=10,
        bold=True,
        color=TEAL,
        font_name=FONT_BODY,
        center=True,
    )
    hero = _add_textbox(s1, 7680960, 5577840, 4297680, 700000)
    headline = _as_list(t1.get("hero_headline"))[:3] or [_truncate(str(result.get("recommended_basis") or ""), 36)]
    _fill_textbox(
        hero,
        [_truncate(line, 32) for line in headline],
        width_emu=4297680,
        height_emu=700000,
        preferred_pt=14,
        min_pt=10,
        bold=True,
        color=NAVY,
        font_name=FONT_BODY,
        center=True,
    )

    # --- Slide 2: Design basis cards ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    t2 = slides[1] if isinstance(slides[1], Mapping) else {}
    _eyebrow_and_title(
        s2,
        str(t2.get("eyebrow") or f"Agitator Selection / {system}"),
        str(t2.get("heading") or "Design basis from DIR code"),
    )
    cards = [c for c in (t2.get("cards") or []) if isinstance(c, Mapping)][:6]
    positions = [
        (658368, 1508760),
        (4389120, 1508760),
        (8092440, 1508760),
        (658368, 3337560),
        (4389120, 3337560),
        (8092440, 3337560),
    ]
    for (left, top), card in zip(positions, cards):
        _add_rect(s2, left, top, 3246120, 1216152, CARD_FILL)
        lab = _add_textbox(s2, left + 201168, top + 164592, 2834640, 201168)
        _set_run(
            lab.text_frame.paragraphs[0],
            _truncate(card.get("label") or "", 28),
            size_pt=9.5,
            bold=True,
            color=GRAY,
            font_name=FONT_BODY,
        )
        val = _add_textbox(s2, left + 201168, top + 493776, 2788920, 500000)
        accent = bool(card.get("accent"))
        _fill_textbox(
            val,
            _truncate(card.get("value") or "", 56),
            width_emu=2788920,
            height_emu=500000,
            preferred_pt=14,
            min_pt=10,
            bold=True,
            color=TEAL if accent else NAVY,
            font_name=FONT_BODY,
        )
    impl_l = _add_textbox(s2, 658368, 5102352, 2194560, 274320)
    _set_run(impl_l.text_frame.paragraphs[0], "Selection implication", size_pt=15, bold=True, color=NAVY, font_name=FONT_BODY)
    impl = _add_textbox(s2, 2971800, 5056632, 8458200, 700000)
    _fill_textbox(
        impl,
        _truncate(t2.get("selection_implication") or "", 200),
        width_emu=8458200,
        height_emu=700000,
        preferred_pt=14,
        min_pt=10,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
    )
    _footer(s2, slide_no=2, dir_code=dir_code)

    # --- Slide 3: Objectives / failure modes ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    t3 = slides[2] if isinstance(slides[2], Mapping) else {}
    _eyebrow_and_title(
        s3,
        str(t3.get("eyebrow") or f"Agitator Selection / {system}"),
        str(t3.get("heading") or "Mixing objectives, constraints and failure modes"),
    )
    steps = [x for x in (t3.get("process_steps") or []) if isinstance(x, Mapping)][:4]
    step_lefts = [658368, 3218688, 5779008, 8339328]
    for left, step in zip(step_lefts, steps):
        _add_rect(s3, left, 2057400, 1965960, 1143000, CARD_FILL)
        # number circle (approx with rounded rect)
        _add_rect(s3, left + 109728, 2258568, 329184, 329184, TEAL)
        num = _add_textbox(s3, left + 219456, 2331720, 109728, 109728)
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_run(num.text_frame.paragraphs[0], str(step.get("n") or ""), size_pt=10, bold=True, color=WHITE, font_name=FONT_BODY)
        title = _add_textbox(s3, left + 530352, 2240280, 1325880, 219456)
        _set_run(
            title.text_frame.paragraphs[0],
            _truncate(step.get("title") or "", 22),
            size_pt=14,
            bold=True,
            color=NAVY,
            font_name=FONT_BODY,
        )
        detail = _add_textbox(s3, left + 530352, 2587752, 1325880, 480000)
        _fill_textbox(
            detail,
            _truncate(step.get("detail") or "", 48),
            width_emu=1325880,
            height_emu=480000,
            preferred_pt=11,
            min_pt=9,
            bold=False,
            color=GRAY,
            font_name=FONT_BODY,
        )
    h_fail = _add_textbox(s3, 658368, 4069080, 2743200, 320040)
    _set_run(h_fail.text_frame.paragraphs[0], "What can go wrong", size_pt=16, bold=True, color=NAVY, font_name=FONT_BODY)
    fail_box = _add_textbox(s3, 777240, 4590288, 5394960, 1280160)
    _fill_textbox(
        fail_box,
        _as_list(t3.get("failure_modes")),
        width_emu=5394960,
        height_emu=1280160,
        preferred_pt=12,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
        bullet=True,
        max_items=4,
        char_cap=78,
    )
    h_out = _add_textbox(s3, 6446520, 4069080, 2743200, 320040)
    _set_run(h_out.text_frame.paragraphs[0], "Target outcome", size_pt=16, bold=True, color=NAVY, font_name=FONT_BODY)
    _add_rect(s3, 6492240, 4535424, 4773168, 1234440, CHIP_TEAL_FILL)
    out_box = _add_textbox(s3, 6812280, 4828032, 4160520, 800000)
    _fill_textbox(
        out_box,
        _truncate(t3.get("target_outcome") or "", 140),
        width_emu=4160520,
        height_emu=800000,
        preferred_pt=13,
        min_pt=10,
        bold=True,
        color=NAVY,
        font_name=FONT_BODY,
        center=True,
    )
    _footer(s3, slide_no=3, dir_code=dir_code)

    # --- Slide 4: Options shortlist ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    t4 = slides[3] if isinstance(slides[3], Mapping) else {}
    _eyebrow_and_title(
        s4,
        str(t4.get("eyebrow") or f"Agitator Selection / {system}"),
        str(t4.get("heading") or "Realistic mixing-system options"),
    )
    headers = [("Mixing type", 685800), ("Notes", 4434840), ("Fit", 8915400)]
    for label, left in headers:
        box = _add_textbox(s4, left, 1417320, 2000000 if left != 685800 else 3500000, 237744)
        _set_run(box.text_frame.paragraphs[0], label, size_pt=12, bold=True, color=GRAY, font_name=FONT_BODY)
    _add_line(s4, 685800, 1764792, 9509760)
    rows = [r for r in (t4.get("rows") or []) if isinstance(r, Mapping)][:6]
    y = 2011680
    for row in rows:
        name = _add_textbox(s4, 685800, y, 3600000, 360000)
        fit = str(row.get("fit") or "").lower()
        name_color = TEAL if fit in {"best", "recommended"} else BODY
        _fill_textbox(
            name,
            _truncate(row.get("name") or "", 48),
            width_emu=3600000,
            height_emu=360000,
            preferred_pt=12,
            min_pt=9,
            bold=fit in {"best", "recommended"},
            color=name_color,
            font_name=FONT_BODY,
        )
        notes = _add_textbox(s4, 4434840, y, 4200000, 360000)
        _fill_textbox(
            notes,
            _truncate(row.get("notes") or "", 70),
            width_emu=4200000,
            height_emu=360000,
            preferred_pt=11,
            min_pt=9,
            bold=False,
            color=BODY,
            font_name=FONT_BODY,
        )
        fit_box = _add_textbox(s4, 8915400, y, 1645920, 360000)
        _fill_textbox(
            fit_box,
            _truncate(str(row.get("fit") or "").title(), 16),
            width_emu=1645920,
            height_emu=360000,
            preferred_pt=12,
            min_pt=9,
            bold=fit in {"best", "recommended"},
            color=TEAL if fit in {"best", "recommended"} else BODY,
            font_name=FONT_BODY,
            center=True,
        )
        y += 520000
    rec = _add_textbox(s4, 685800, 5600000, 10400000, 600000)
    _fill_textbox(
        rec,
        _truncate(t4.get("recommendation_line") or "", 180),
        width_emu=10400000,
        height_emu=600000,
        preferred_pt=13,
        min_pt=10,
        bold=True,
        color=NAVY,
        font_name=FONT_BODY,
    )
    _footer(s4, slide_no=4, dir_code=dir_code)

    # --- Slide 5: Matrix ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    t5 = slides[4] if isinstance(slides[4], Mapping) else {}
    _eyebrow_and_title(
        s5,
        str(t5.get("eyebrow") or f"Agitator Selection / {system}"),
        str(t5.get("heading") or "Option evaluation matrix"),
    )
    cols = [
        ("Option", 685800, 2800000),
        ("Technical", 3600000, 1400000),
        ("GMP", 5100000, 1000000),
        ("Scale-up", 6200000, 1200000),
        ("Cost", 7500000, 1200000),
        ("Reliability", 8800000, 1200000),
        ("Rank", 10100000, 800000),
    ]
    for label, left, width in cols:
        box = _add_textbox(s5, left, 1417320, width, 237744)
        _set_run(box.text_frame.paragraphs[0], label, size_pt=11, bold=True, color=GRAY, font_name=FONT_BODY)
    _add_line(s5, 685800, 1764792, 10400000)
    mrows = [r for r in (t5.get("rows") or []) if isinstance(r, Mapping)][:5]
    y = 2011680
    for row in mrows:
        values = [
            (row.get("option"), 685800, 2800000),
            (row.get("technical_fit"), 3600000, 1400000),
            (row.get("gmp"), 5100000, 1000000),
            (row.get("scale_up_risk"), 6200000, 1200000),
            (row.get("cost_schedule"), 7500000, 1200000),
            (row.get("reliability"), 8800000, 1200000),
            (row.get("rank"), 10100000, 800000),
        ]
        for i, (val, left, width) in enumerate(values):
            box = _add_textbox(s5, left, y, width, 340000)
            _fill_textbox(
                box,
                _truncate(str(val or ""), 40 if i == 0 else 18),
                width_emu=width,
                height_emu=340000,
                preferred_pt=11,
                min_pt=9,
                bold=(i == 0 and str(row.get("rank") or "") in {"1", "1.0"}),
                color=TEAL if (i == 0 and str(row.get("rank") or "") in {"1", "1.0"}) else BODY,
                font_name=FONT_BODY,
            )
        y += 480000
    logic = _add_textbox(s5, 685800, 5000000, 10400000, 1000000)
    _fill_textbox(
        logic,
        "Decision logic  " + _truncate(t5.get("decision_logic") or "", 200),
        width_emu=10400000,
        height_emu=1000000,
        preferred_pt=12,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
    )
    _footer(s5, slide_no=5, dir_code=dir_code)

    # --- Slide 6: Recommendation ---
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    t6 = slides[5] if isinstance(slides[5], Mapping) else {}
    _eyebrow_and_title(
        s6,
        str(t6.get("eyebrow") or f"Agitator Selection / {system}"),
        str(t6.get("heading") or "Recommended basis and alternate option"),
    )
    _add_rect(s6, 658368, 1400000, 5600000, 4200000, CARD_FILL)
    h = _add_textbox(s6, 900000, 1600000, 5000000, 400000)
    _set_run(h.text_frame.paragraphs[0], "Recommended selection", size_pt=13, bold=True, color=GRAY, font_name=FONT_BODY)
    rec = _add_textbox(s6, 900000, 2000000, 5000000, 900000)
    _fill_textbox(
        rec,
        _truncate(t6.get("recommended") or "", 110),
        width_emu=5000000,
        height_emu=900000,
        preferred_pt=16,
        min_pt=11,
        bold=True,
        color=TEAL,
        font_name=FONT_DISPLAY,
    )
    why = _add_textbox(s6, 900000, 3000000, 5000000, 2300000)
    _fill_textbox(
        why,
        _as_list(t6.get("recommended_why")) or _as_list(t6.get("pros")),
        width_emu=5000000,
        height_emu=2300000,
        preferred_pt=12,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
        bullet=True,
        max_items=4,
        char_cap=72,
    )
    _add_rect(s6, 6600000, 1400000, 5000000, 2000000, CHIP_TEAL_FILL)
    alt_h = _add_textbox(s6, 6850000, 1550000, 4500000, 300000)
    _set_run(alt_h.text_frame.paragraphs[0], "Alternate option", size_pt=13, bold=True, color=GRAY, font_name=FONT_BODY)
    alt = _add_textbox(s6, 6850000, 1900000, 4500000, 1200000)
    _fill_textbox(
        alt,
        _truncate(t6.get("alternate") or "", 100),
        width_emu=4500000,
        height_emu=1200000,
        preferred_pt=14,
        min_pt=10,
        bold=True,
        color=NAVY,
        font_name=FONT_BODY,
    )
    alt_n = _add_textbox(s6, 6850000, 3000000, 4500000, 300000)
    _fill_textbox(
        alt_n,
        _truncate(t6.get("alternate_note") or "", 100),
        width_emu=4500000,
        height_emu=300000,
        preferred_pt=11,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
    )
    cons_h = _add_textbox(s6, 6600000, 3600000, 5000000, 300000)
    _set_run(cons_h.text_frame.paragraphs[0], "Cons / watchouts", size_pt=13, bold=True, color=GRAY, font_name=FONT_BODY)
    cons = _add_textbox(s6, 6600000, 3950000, 5000000, 1400000)
    _fill_textbox(
        cons,
        _as_list(t6.get("cons")),
        width_emu=5000000,
        height_emu=1400000,
        preferred_pt=11,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
        bullet=True,
        max_items=4,
        char_cap=64,
    )
    _footer(s6, slide_no=6, dir_code=dir_code)

    # --- Slide 7: Specs / vendors ---
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    t7 = slides[6] if isinstance(slides[6], Mapping) else {}
    _eyebrow_and_title(
        s7,
        str(t7.get("eyebrow") or f"Agitator Selection / {system}"),
        str(t7.get("heading") or "Preliminary specification points / vendors / references"),
    )
    _add_rect(s7, 658368, 1400000, 5600000, 4500000, CARD_FILL)
    sh = _add_textbox(s7, 900000, 1550000, 5000000, 300000)
    _set_run(sh.text_frame.paragraphs[0], "Baseline specification", size_pt=14, bold=True, color=NAVY, font_name=FONT_BODY)
    specs = _add_textbox(s7, 900000, 2000000, 5000000, 3600000)
    _fill_textbox(
        specs,
        _as_list(t7.get("specs")),
        width_emu=5000000,
        height_emu=3600000,
        preferred_pt=12,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
        bullet=True,
        max_items=7,
        char_cap=72,
    )
    mh = _add_textbox(s7, 6600000, 1400000, 5000000, 300000)
    _set_run(mh.text_frame.paragraphs[0], "Manufacturer shortlist", size_pt=14, bold=True, color=NAVY, font_name=FONT_BODY)
    mfrs = _add_textbox(s7, 6600000, 1800000, 5000000, 1800000)
    _fill_textbox(
        mfrs,
        _as_list(t7.get("manufacturers")),
        width_emu=5000000,
        height_emu=1800000,
        preferred_pt=12,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
        bullet=True,
        max_items=6,
        char_cap=60,
    )
    dh = _add_textbox(s7, 6600000, 3800000, 5000000, 300000)
    _set_run(dh.text_frame.paragraphs[0], "Do not specify", size_pt=14, bold=True, color=NAVY, font_name=FONT_BODY)
    dns = _add_textbox(s7, 6600000, 4200000, 5000000, 1400000)
    _fill_textbox(
        dns,
        _as_list(t7.get("do_not_specify")),
        width_emu=5000000,
        height_emu=1400000,
        preferred_pt=11,
        min_pt=9,
        bold=False,
        color=BODY,
        font_name=FONT_BODY,
        bullet=True,
        max_items=5,
        char_cap=58,
    )
    _footer(s7, slide_no=7, dir_code=dir_code)

    # Save with lock-safe fallback (Windows PermissionError if PPTX is open).
    try:
        prs.save(str(out))
        return out
    except PermissionError:
        from datetime import datetime

        stamped = out.with_name(f"{out.stem}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{out.suffix}")
        prs.save(str(stamped))
        return stamped


# Back-compat alias used by tests/imports
def build_evaluation_pptx_simple(*args, **kwargs):
    return build_evaluation_pptx(*args, **kwargs)
